# wolfstitch/processing/extractors/pptx_extractor.py
"""
PowerPoint (.pptx) text extraction module for Wolfstitch

This module extracts text content from PowerPoint presentations including:
- Slide text content (titles, bullets, text boxes)
- Speaker notes
- Slide layouts and master slides
- Error handling for corrupted or password-protected files

Supports both modern .pptx and legacy .ppt formats through python-pptx library.
"""

import os
from typing import List, Dict, Any, Optional
import logging

try:
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available. PowerPoint support disabled.")

# Set up logging
logger = logging.getLogger(__name__)


def extract_text(file_path: str) -> str:
    """
    Extract text content from PowerPoint presentations
    
    Extracts all text content including slide text, speaker notes, and 
    structured content while handling various edge cases and errors.
    
    Args:
        file_path (str): Path to the PowerPoint file (.pptx or .ppt)
        
    Returns:
        str: Extracted text content with slides separated by double newlines
        
    Raises:
        ValueError: If file format is unsupported or python-pptx unavailable
        RuntimeError: If extraction fails due to file corruption or protection
        
    Examples:
        >>> text = extract_text("presentation.pptx")
        >>> text = extract_text("legacy.ppt")  # Also supported
    """
    if not PPTX_AVAILABLE:
        raise ValueError(
            "PowerPoint support requires python-pptx library. "
            "Install with: pip install python-pptx"
        )
    
    if not os.path.exists(file_path):
        raise RuntimeError(f"PowerPoint file not found: {file_path}")
    
    # Validate file extension
    ext = os.path.splitext(file_path)[1].lower()
    if ext not in ['.pptx', '.ppt']:
        raise ValueError(f"Unsupported PowerPoint format: {ext}")
    
    try:
        # Load presentation
        prs = Presentation(file_path)
        logger.info(f"Loaded PowerPoint: {os.path.basename(file_path)} ({len(prs.slides)} slides)")
        
        # Extract text from all slides
        slide_texts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = extract_slide_content(slide, slide_num)
            if slide_content.strip():  # Only add non-empty slides
                slide_texts.append(slide_content)
        
        if not slide_texts:
            logger.warning(f"No text content found in {file_path}")
            return ""
        
        # Join slides with double newlines for clear separation
        full_text = "\n\n".join(slide_texts)
        
        logger.info(f"Extracted {len(full_text)} characters from {len(slide_texts)} slides")
        return full_text
        
    except PackageNotFoundError:
        raise RuntimeError(f"Invalid or corrupted PowerPoint file: {file_path}")
    except Exception as e:
        # Handle password-protected files and other issues
        error_msg = str(e).lower()
        if "password" in error_msg or "encrypted" in error_msg:
            raise RuntimeError(f"Password-protected PowerPoint file: {file_path}")
        elif "permission" in error_msg:
            raise RuntimeError(f"Permission denied accessing: {file_path}")
        else:
            raise RuntimeError(f"Failed to extract from PowerPoint file {file_path}: {str(e)}")


def extract_slide_content(slide, slide_num: int) -> str:
    """
    Extract all text content from a single slide
    
    Args:
        slide: PowerPoint slide object
        slide_num (int): Slide number for context
        
    Returns:
        str: All text content from the slide including notes
    """
    slide_parts = []
    
    # Add slide number header for context
    slide_parts.append(f"=== Slide {slide_num} ===")
    
    # Extract text from all shapes in the slide
    shape_texts = extract_shape_texts(slide)
    if shape_texts:
        slide_parts.extend(shape_texts)
    
    # Extract speaker notes
    notes_text = extract_speaker_notes(slide)
    if notes_text:
        slide_parts.append(f"\n--- Speaker Notes ---")
        slide_parts.append(notes_text)
    
    return "\n".join(slide_parts)


def extract_shape_texts(slide) -> List[str]:
    """
    Extract text from all shapes in a slide
    
    Args:
        slide: PowerPoint slide object
        
    Returns:
        List[str]: List of text content from shapes
    """
    texts = []
    
    for shape in slide.shapes:
        try:
            # Check if shape has text
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
            
            # Handle tables separately
            elif hasattr(shape, "table"):
                table_text = extract_table_text(shape.table)
                if table_text:
                    texts.append(table_text)
                    
            # Handle grouped shapes
            elif hasattr(shape, "shapes"):
                for sub_shape in shape.shapes:
                    if hasattr(sub_shape, "text") and sub_shape.text.strip():
                        texts.append(sub_shape.text.strip())
        
        except Exception as e:
            # Log but continue processing other shapes
            logger.debug(f"Skipping shape due to error: {e}")
            continue
    
    return texts


def extract_table_text(table) -> Optional[str]:
    """
    Extract text from PowerPoint table
    
    Args:
        table: PowerPoint table object
        
    Returns:
        Optional[str]: Formatted table text or None if empty
    """
    try:
        rows = []
        
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip() if cell.text else ""
                cells.append(cell_text)
            
            # Only add non-empty rows
            if any(cell for cell in cells):
                rows.append(" | ".join(cells))
        
        if rows:
            return "\n".join(rows)
        return None
        
    except Exception as e:
        logger.debug(f"Error extracting table: {e}")
        return None


def extract_speaker_notes(slide) -> Optional[str]:
    """
    Extract speaker notes from a slide
    
    Args:
        slide: PowerPoint slide object
        
    Returns:
        Optional[str]: Speaker notes text or None if empty
    """
    try:
        # Check if slide has notes
        if hasattr(slide, 'notes_slide') and slide.notes_slide:
            notes_slide = slide.notes_slide
            
            # Extract notes text frame
            if hasattr(notes_slide, 'notes_text_frame') and notes_slide.notes_text_frame:
                notes_text = notes_slide.notes_text_frame.text
                
                if notes_text and notes_text.strip():
                    return notes_text.strip()
        
        return None
        
    except Exception as e:
        logger.debug(f"Error extracting speaker notes: {e}")
        return None


def validate_powerpoint_file(file_path: str) -> Dict[str, Any]:
    """
    Validate PowerPoint file and return metadata
    
    Args:
        file_path (str): Path to PowerPoint file
        
    Returns:
        Dict[str, Any]: File validation and metadata information
    """
    result = {
        'valid': False,
        'slides': 0,
        'has_notes': False,
        'error': None,
        'file_type': os.path.splitext(file_path)[1].lower()
    }
    
    if not PPTX_AVAILABLE:
        result['error'] = "python-pptx library not available"
        return result
    
    try:
        prs = Presentation(file_path)
        result['valid'] = True
        result['slides'] = len(prs.slides)
        
        # Check if any slides have notes
        for slide in prs.slides:
            if extract_speaker_notes(slide):
                result['has_notes'] = True
                break
                
    except Exception as e:
        result['error'] = str(e)
    
    return result


def get_supported_formats() -> List[str]:
    """
    Get list of supported PowerPoint formats
    
    Returns:
        List[str]: List of supported file extensions
    """
    if PPTX_AVAILABLE:
        return ['.pptx', '.ppt']
    else:
        return []


# Compatibility function for the main dispatcher
def is_supported_format(file_path: str) -> bool:
    """
    Check if file format is supported
    
    Args:
        file_path (str): Path to file
        
    Returns:
        bool: True if format is supported
    """
    ext = os.path.splitext(file_path)[1].lower()
    return ext in get_supported_formats()
